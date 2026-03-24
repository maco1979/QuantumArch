from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

PRIMARY = RGBColor(0x1A, 0x5F, 0x7A)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
RED = RGBColor(0xF8, 0x71, 0x71)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
PINK = RGBColor(0xF4, 0x72, 0xB6)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    p = sub_box.text_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_chart_title(slide, title):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.alignment = PP_ALIGN.CENTER

# 1. 封面
add_title_slide(prs, '熊猫讲技术分析', '可视化示例图解 - 8大核心图表')

# 2. 胜率盈亏比对比图（柱状图）
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图1：胜率与盈亏比对比（期望值）')

chart_data = ChartData()
chart_data.categories = ['高频小赚', '趋势跟踪', '平衡策略']
chart_data.add_series('胜率', (0.70, 0.40, 0.50))
chart_data.add_series('盈亏比', (0.80, 3.00, 2.00))
chart_data.add_series('期望值', (0.56, 1.20, 1.00))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(11.333), Inches(5),
    chart_data
).chart

box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.8))
p = box.text_frame.paragraphs[0]
p.text = '结论：40%胜率x3:1盈亏比(期望值1.2) > 70%胜率x0.8:1盈亏比(期望值0.56)'
p.font.size = Pt(16)
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER

# 3. MACD战法示意图
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图2：MACD只吃一口战法 - 红绿柱变化')

chart_data = ChartData()
chart_data.categories = ['1', '2', '3', '4', '5', '6', '7', '8', '9', '10', '11', '12']
macd_values = [-80, -70, -60, -40, -20, -10, -5, 10, 30, 50, 60, 40]
chart_data.add_series('MACD', macd_values)

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5),
    chart_data
).chart

box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(1))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '口诀：买在小绿柱（下跌力度衰竭）→ 卖在小红柱（上涨动力不足）'
p.font.size = Pt(18)
p.font.color.rgb = YELLOW
p.alignment = PP_ALIGN.CENTER

# 4. 复盘账户增长对比（折线图）
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图3：复盘 vs 不复盘 - 账户增长曲线对比')

chart_data = ChartData()
chart_data.categories = ['第1月', '第2月', '第3月', '第4月', '第5月', '第6月', '第7月', '第8月', '第9月', '第10月', '第11月', '第12月']
chart_data.add_series('不复盘（震荡下行）', (100, 95, 98, 90, 92, 85, 88, 80, 75, 70, 65, 60))
chart_data.add_series('认真复盘（复利增长）', (100, 102, 105, 108, 112, 118, 125, 135, 148, 165, 185, 210))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5),
    chart_data
).chart

box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
p = box.text_frame.paragraphs[0]
p.text = '复盘的质量，决定了你账户的厚度！'
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = GREEN
p.alignment = PP_ALIGN.CENTER

# 5. 亏损原因饼图
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图4：亏损原因占比分析（饼图）')

chart_data = ChartData()
chart_data.categories = ['没有交易系统', '不执行纪律', '仓位管理不当', '心态不稳定', '不复盘总结']
chart_data.add_series('占比', (35, 25, 20, 15, 5))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(2.5), Inches(1.5), Inches(8.333), Inches(5),
    chart_data
).chart

box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.8))
p = box.text_frame.paragraphs[0]
p.text = '最大问题：没有交易系统(35%) + 不执行纪律(25%) = 60%'
p.font.size = Pt(16)
p.font.color.rgb = RED
p.alignment = PP_ALIGN.CENTER

# 6. 学习路径阶梯图（柱状图）
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图5：交易学习路径 - 五阶段能力成长')

chart_data = ChartData()
chart_data.categories = ['基础概念', '技术指标', '交易系统', '资金管理', '心态修炼']
chart_data.add_series('能力值(%)', (20, 40, 60, 75, 90))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5),
    chart_data
).chart

box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
p = box.text_frame.paragraphs[0]
p.text = '循序渐进：基础 → 指标 → 系统 → 资金 → 心态'
p.font.size = Pt(18)
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER

# 7. 过滤漏斗（柱状图）
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图6：交易机会过滤漏斗')

chart_data = ChartData()
chart_data.categories = ['所有机会', '趋势明确', '信号清晰', '盈亏比合理', '最终执行']
chart_data.add_series('剩余比例(%)', (100, 60, 35, 20, 10))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5),
    chart_data
).chart

box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
p = box.text_frame.paragraphs[0]
p.text = '过滤掉90%的低质量机会，只抓最好的10%！'
p.font.size = Pt(18)
p.font.color.rgb = PURPLE
p.alignment = PP_ALIGN.CENTER

# 8. 六要素示意图
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图7：交易系统六要素闭环')

shapes = slide.shapes
center_x, center_y = Inches(6.666), Inches(4.2)
radius = Inches(2.2)

elements = [
    ('方向', 90, BLUE), ('位置', 30, PURPLE), ('时机', -30, PINK),
    ('仓位', -90, YELLOW), ('止损', -150, RED), ('止盈', 150, GREEN)
]

import math
for elem, angle, color in elements:
    rad = math.radians(angle)
    x = center_x + radius * math.cos(rad)
    y = center_y - radius * math.sin(rad)
    
    circle = shapes.add_shape(MSO_SHAPE.OVAL, x - Inches(0.5), y - Inches(0.5), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    
    box = shapes.add_textbox(x - Inches(0.4), y - Inches(0.2), Inches(0.8), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = elem
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

circle = shapes.add_shape(MSO_SHAPE.OVAL, center_x - Inches(0.7), center_y - Inches(0.7), Inches(1.4), Inches(1.4))
circle.fill.solid()
circle.fill.fore_color.rgb = PRIMARY

box = shapes.add_textbox(center_x - Inches(0.6), center_y - Inches(0.3), Inches(1.2), Inches(0.6))
p = box.text_frame.paragraphs[0]
p.text = '交易系统'
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.8))
p = box.text_frame.paragraphs[0]
p.text = '六要素缺一不可：方向→位置→时机→仓位→止损→止盈'
p.font.size = Pt(16)
p.font.color.rgb = PRIMARY
p.alignment = PP_ALIGN.CENTER

# 9. 缠论概念（折线图）
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_chart_title(slide, '图8：缠论核心 - 中枢、背驰与买卖点')

chart_data = ChartData()
chart_data.categories = ['1', '2', '3', '4', '5', '6', '7', '8', '9', '10', '11', '12', '13', '14', '15']
chart_data.add_series('价格走势', (100, 120, 110, 130, 125, 115, 105, 95, 100, 110, 125, 140, 135, 120, 100))
chart_data.add_series('MACD', (-30, -20, -10, 20, 10, -5, -25, -35, -20, 10, 40, 60, 40, 10, -30))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE, Inches(1), Inches(1.5), Inches(11.333), Inches(4.5),
    chart_data
).chart

box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.333), Inches(0.8))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '底背驰（价格新低+MACD绿柱缩短）= 买点 | 顶背驰（价格新高+MACD红柱缩短）= 卖点'
p.font.size = Pt(16)
p.font.color.rgb = YELLOW
p.alignment = PP_ALIGN.CENTER

# 10. 总结
slide = prs.slides.add_slide(prs.slide_layouts[6])
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = PRIMARY

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1))
p = box.text_frame.paragraphs[0]
p.text = '交易成功的关键'
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10.333), Inches(3.5))
tf = box.text_frame
tf.word_wrap = True

points = [
    ('技术层面：', True, YELLOW),
    ('  胜率x盈亏比>1是盈利基础', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('  MACD只吃一口战法简单有效', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('  六要素框架构建完整系统', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('执行层面：', True, YELLOW),
    ('  严格资金管理（2%原则）', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('  坚持每日复盘（持续改进）', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('  知道不做什么（过滤机会）', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('', False, RGBColor(0xFF, 0xFF, 0xFF)),
    ('记住：交易是概率游戏，不是赌博！', True, YELLOW)
]

for i, (text, bold, color) in enumerate(points):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(20)
    p.font.bold = bold
    p.font.color.rgb = color
    p.space_before = Pt(8)

prs.save('熊猫讲技术分析_可视化图表.pptx')
print('可视化图表PPT创建成功！')
